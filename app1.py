
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10) # Standard widescreen 16:9
    prs.slide_height = Inches(5.625)

    # --- Common Styles ---
    TEXT_COLOR = RGBColor(50, 50, 50) # Dark Grey
    TITLE_COLOR = RGBColor(0, 51, 102) # Dark Blue

    # Function to add a slide with a title and content placeholder
    def add_content_slide(prs, title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Title placeholder
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.name = 'Calibri Light'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Body placeholder
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear existing paragraphs

        for point_data in bullet_points:
            p = tf.add_paragraph()
            p.text = point_data['text']
            p.level = point_data['level']
            p.font.name = 'Calibri'
            p.font.size = Pt(22 - (point_data['level'] * 2)) # Adjust size for levels (22pt, 20pt, 18pt)
            p.font.color.rgb = TEXT_COLOR
            tf.word_wrap = True

        # Adjust body placeholder position for better margins and consistent look
        body.top = Inches(1.5)
        body.left = Inches(0.7)
        body.width = Inches(8.6)
        body.height = Inches(3.8) # Allow enough height for bullets

    # --- Slide 1: Title Slide (Overall Presentation Title) ---
    slide_layout = prs.slide_layouts[0] # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Data Science with Gen AI and Agentic AI"
    title.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = slide.placeholders[1]
    subtitle.text = "Exploring the Next Frontier in Data-Driven Innovation"
    subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add a date/presenter text box manually for better control and fixed position
    left = Inches(1.0)
    top = Inches(4.5)
    width = Inches(8.0)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Presented by [Your Name] | [Date]" # Placeholder: Customize your name and date
    p.font.name = 'Calibri'
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Introduction to Data Science ---
    add_content_slide(prs, "Introduction to Data Science", [
        {'text': 'Interdisciplinary field combining statistics, computer science, and domain expertise.', 'level': 0},
        {'text': 'Focuses on extracting knowledge and insights from structured and unstructured data.', 'level': 0},
        {'text': 'Key phases: Data Collection, Cleaning, Analysis, Modeling, and Visualization.', 'level': 0},
        {'text': 'Goal: Drive informed decision-making and foster innovation across industries.', 'level': 0}
    ])

    # --- Slide 3: Introduction to Generative AI ---
    add_content_slide(prs, "Introduction to Generative AI", [
        {'text': 'Branch of AI focused on creating new, original content and data.', 'level': 0},
        {'text': 'Learns complex patterns and structures from large datasets.', 'level': 0},
        {'text': 'Capable of generating diverse outputs: text, images, audio, video, and code.', 'level': 0},
        {'text': 'Examples include Large Language Models (LLMs) like GPT and image models like DALL-E.', 'level': 0}
    ])

    # --- Slide 4: Introduction to Agentic AI ---
    add_content_slide(prs, "Introduction to Agentic AI", [
        {'text': 'AI systems designed to operate autonomously, pursuing defined goals.', 'level': 0},
        {'text': 'Characterized by perception, planning, action execution, and memory.', 'level': 0},
        {'text': 'Can break down complex tasks into manageable sub-tasks and execute them sequentially.', 'level': 0},
        {'text': 'Often leverage LLMs for reasoning, decision-making, and interaction with tools.', 'level': 0}
    ])

    # --- Slide 5: How Gen AI Enhances Data Science ---
    add_content_slide(prs, "How Gen AI Enhances Data Science", [
        {'text': 'Data Augmentation:', 'level': 0},
        {'text': 'Synthesizing realistic training data, especially for rare or sensitive cases.', 'level': 1},
        {'text': 'Automated Feature Engineering:', 'level': 0},
        {'text': 'Discovering and creating new, relevant features from raw data.', 'level': 1},
        {'text': 'Model Generation & Optimization:', 'level': 0},
        {'text': 'Assisting in architecture design and hyperparameter tuning for machine learning models.', 'level': 1},
        {'text': 'Code Generation:', 'level': 0},
        {'text': 'Automating script writing for data cleaning, analysis, and ML pipeline development.', 'level': 1},
        {'text': 'Insight Generation:', 'level': 0},
        {'text': 'Summarizing complex findings and generating natural language reports.', 'level': 1}
    ])

    # --- Slide 6: Architecture of Agentic AI Systems ---
    add_content_slide(prs, "Architecture of Agentic AI Systems", [
        {'text': 'LLM Core (The Brain):', 'level': 0},
        {'text': 'Provides reasoning, planning capabilities, and guides overall behavior.', 'level': 1},
        {'text': 'Memory Module:', 'level': 0},
        {'text': 'Stores past interactions, observations, and long-term knowledge for context.', 'level': 1},
        {'text': 'Perception Module:', 'level': 0},
        {'text': 'Gathers information from the environment (e.g., databases, APIs, web).', 'level': 1},
        {'text': 'Planning Module:', 'level': 0},
        {'text': 'Decomposes high-level goals into actionable, executable steps.', 'level': 1},
        {'text': 'Tool Use / Action Module:', 'level': 0},
        {'text': 'Accesses and utilizes external tools (e.g., code interpreters, web search, APIs) to perform tasks.', 'level': 1}
    ])

    # --- Slide 7: Agentic AI vs Gen AI ---
    add_content_slide(prs, "Agentic AI vs Gen AI", [
        {'text': 'Generative AI:', 'level': 0},
        {'text': 'Focus: Creation of novel data or content based on learned patterns.', 'level': 1},
        {'text': 'Primary Goal: Generate output that resembles or extends training data.', 'level': 1},
        {'text': 'Key Capability: Pattern recognition, synthesis, and mimicry.', 'level': 1},
        {'text': 'Autonomy: Limited; typically responds to specific prompts or inputs.', 'level': 1},
        {'text': 'Agentic AI:', 'level': 0},
        {'text': 'Focus: Autonomous, goal-oriented action and complex problem-solving.', 'level': 1},
        {'text': 'Primary Goal: Achieve predefined objectives through planning and execution.', 'level': 1},
        {'text': 'Key Capability: Reasoning, planning, tool orchestration, and self-correction.', 'level': 1},
        {'text': 'Autonomy: High; operates independently and adaptively towards a goal.', 'level': 1}
    ])

    # --- Slide 8: Real-world Applications of Gen AI and Agentic AI ---
    add_content_slide(prs, "Real-world Applications of Gen AI and Agentic AI", [
        {'text': 'Generative AI Applications:', 'level': 0},
        {'text': 'Personalized content creation (marketing, news articles).', 'level': 1},
        {'text': 'Accelerated drug discovery and material science design.', 'level': 1},
        {'text': 'Synthetic data generation for privacy-preserving research and testing.', 'level': 1},
        {'text': 'Art, music, and story creation.', 'level': 1},
        {'text': 'Agentic AI Applications:', 'level': 0},
        {'text': 'Autonomous customer support and virtual assistants.', 'level': 1},
        {'text': 'Automated software development, testing, and debugging.', 'level': 1},
        {'text': 'Personalized research and information gathering agents.', 'level': 1},
        {'text': 'Complex task automation across various business processes.', 'level': 1}
    ])

    # --- Slide 9: Future of Data Science with Gen AI and Agentic AI ---
    add_content_slide(prs, "Future of Data Science with Gen AI and Agentic AI", [
        {'text': 'Enhanced Automation:', 'level': 0},
        {'text': 'Automating data pipelines, model building, and insight generation from raw data.', 'level': 1},
        {'text': 'Democratization of Data Science:', 'level': 0},
        {'text': 'Making advanced analytics and ML accessible to a broader audience, reducing barriers.', 'level': 1},
        {'text': 'Shift in Roles:', 'level': 0},
        {'text': 'Data scientists evolve to focus more on strategic problem-solving, ethical considerations, and AI governance.', 'level': 1},
        {'text': 'Continuous Learning Systems:', 'level': 0},
        {'text': 'Agentic systems that constantly adapt, learn, and improve data science workflows autonomously.', 'level': 1},
        {'text': 'New Challenges:', 'level': 0},
        {'text': 'Emphasis on ensuring AI safety, mitigating bias, and maintaining interpretability.', 'level': 1}
    ])

    # --- Slide 10: Any Questions? / Thank You ---
    slide_layout = prs.slide_layouts[0] # Title Slide layout for a cleaner 'Thank You' slide
    slide = prs.slides.add_slide(slide_layout)

    # Main text for questions/thank you
    title = slide.shapes.title
    title.text = "Thank You!"
    title.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.top = Inches(1.8) 
    title.height = Inches(1.5)

    # Subtitle/contact info
    subtitle = slide.placeholders[1]
    subtitle.text = "Any Questions? Please connect:"
    subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add contact details as a separate paragraph
    contact_p = subtitle.text_frame.add_paragraph()
    contact_p.text = "[Your Email] | [Your Website/LinkedIn]" # Placeholder: Customize your contact info
    contact_p.font.name = 'Calibri'
    contact_p.font.size = Pt(24)
    contact_p.font.color.rgb = TEXT_COLOR
    contact_p.alignment = PP_ALIGN.CENTER
    
    subtitle.top = Inches(3.8) # Position it lower
    subtitle.height = Inches(1.5)
    subtitle.left = Inches(0.5)
    subtitle.width = Inches(9.0)


    # Save the presentation
    prs.save("Data_Science_with_Gen_AI_and_Agentic_AI.pptx")

create_presentation()
